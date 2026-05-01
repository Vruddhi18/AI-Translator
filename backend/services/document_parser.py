import PyPDF2
from pptx import Presentation
import os

def parse_document(file_path: str, ext: str) -> str:
    text = ""
    try:
        if ext == 'pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        elif ext == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
    except Exception as e:
        print(f"Error parsing document {file_path}: {e}")
        raise e
        
    return text
